#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build a co-author overview presentation (.pptx) from this project's figures and
results. Rich by design -- intended as a starting point to trim/edit manually.

Figures are read from runs/current/output/validation/plots/. Output:
  slides/bap_gapfilling_overview.pptx
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
PLOTS = os.path.join(REPO, "runs", "current", "output", "validation", "plots")
OUTDIR = os.path.join(REPO, "slides")
os.makedirs(OUTDIR, exist_ok=True)

# --- palette ---
DARK = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _tf(box, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    return tf


def header(slide, title, kicker=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), SW - Inches(1.0), Inches(0.8))
    _tf(tb, title, 28, WHITE, bold=True)
    if kicker:
        kb = slide.shapes.add_textbox(Inches(0.5), Inches(0.83), SW - Inches(1.0), Inches(0.3))
        _tf(kb, kicker, 12, RGBColor(0xCF, 0xDA, 0xE8))


def bullets(slide, items, left, top, width, height, size=18, gap=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl, txt = (it if isinstance(it, tuple) else (0, it))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl); r.font.color.rgb = DARK if lvl == 0 else GREY
    return tb


def fit_image(slide, name, box_l, box_t, box_w, box_h):
    path = os.path.join(PLOTS, name)
    if not os.path.exists(path):
        tb = slide.shapes.add_textbox(box_l, box_t, box_w, Inches(0.5))
        _tf(tb, f"[missing figure: {name}]", 14, ACCENT, italic=True)
        return
    iw, ih = Image.open(path).size
    ar = iw / ih
    bw, bh = box_w, box_h
    if bw / bh > ar:           # box wider than image -> fit height
        h = bh; w = int(bh * ar)
    else:                       # fit width
        w = bw; h = int(bw / ar)
    l = box_l + (bw - w) // 2
    t = box_t + (bh - h) // 2
    slide.shapes.add_picture(path, l, t, width=w, height=h)


def takeaway(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), SH - Inches(0.7), SW - Inches(1.0), Inches(0.55))
    _tf(tb, "➤ " + text, 15, ACCENT, bold=True)


def slide_blank():
    return prs.slides.add_slide(BLANK)


def fig_slide(title, fig, side_bullets=None, take=None, kicker=None):
    s = slide_blank(); header(s, title, kicker)
    if side_bullets:
        fit_image(s, fig, Inches(0.4), Inches(1.35), Inches(8.0), Inches(5.3))
        bullets(s, side_bullets, Inches(8.5), Inches(1.5), Inches(4.5), Inches(5.2), size=16)
    else:
        fit_image(s, fig, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.4))
    if take:
        takeaway(s, take)
    return s


# ---------------- 1. Title ----------------
s = slide_blank()
band = s.shapes.add_shape(1, 0, 0, SW, SH); band.fill.solid()
band.fill.fore_color.rgb = DARK; band.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.1), SW - Inches(1.8), Inches(2.4))
_tf(tb, "A globally trained machine-learning model for daily benzo[a]pyrene "
        "gap-filling and a virtual monitoring network over Slovakia", 32, WHITE, bold=True)
ab = s.shapes.add_textbox(Inches(0.9), Inches(4.7), SW - Inches(1.8), Inches(1.0))
_tf(ab, "Juraj Beňo · Dušan Štefánik · Jana Matejovičová", 20,
    RGBColor(0xCF, 0xDA, 0xE8), bold=True)
af = s.shapes.add_textbox(Inches(0.9), Inches(5.3), SW - Inches(1.8), Inches(0.8))
_tf(af, "Slovak Hydrometeorological Institute (SHMÚ)  ·  co-author overview draft", 15,
    RGBColor(0xAB, 0xBD, 0xD4))

# ---------------- 2. The problem ----------------
s = slide_blank(); header(s, "The problem")
bullets(s, [
    "Benzo[a]pyrene (B[a]P) — carcinogenic PAH; EU target value 1 ng m⁻³ (annual mean).",
    "Measurement is laborious: few stations, usually every-third-day sampling.",
    (1, "→ long contiguous gaps in daily records,"),
    (1, "→ most of the territory unmonitored."),
    "Daily fields are needed for exposure, compliance and trend analysis.",
    "Goal: (i) fill daily B[a]P gaps and (ii) extend to a continuous virtual network.",
], Inches(0.7), Inches(1.5), Inches(12), Inches(5), size=20, gap=10)
takeaway(s, "One global model for both gap-filling and spatial extension.")

# ---------------- 3. Approach ----------------
s = slide_blank(); header(s, "Approach")
bullets(s, [
    "Single, globally trained gradient-boosting model (XGBoost) on ln(1+B[a]P).",
    "Inputs: meteorology + co-located proxy pollutants + static spatial covariates.",
    "Honest, leakage-aware validation (three CV designs) against simple baselines.",
    "Quantify what each feature group contributes (ablation + TreeSHAP).",
    "Apply to a continuous meteorological series → virtual monitoring network.",
], Inches(0.7), Inches(1.5), Inches(12), Inches(5), size=20, gap=12)

# ---------------- 4. Data ----------------
s = slide_blank(); header(s, "Data & study area")
bullets(s, [
    "Slovak national B[a]P network, projected onto the 2 km ALADIN grid (215×110).",
    "6 408 daily samples, 21 training stations, 2023-06 to 2025-12.",
    "Station typology (training): 2 rural-bg, 6 suburban-bg, 1 suburban-industrial,",
    (1, "5 urban-bg, 7 urban-traffic."),
    "54 grid-mapped sites total → 33 are virtual (never monitored).",
    "Target: daily B[a]P concentration [ng m⁻³].",
], Inches(0.7), Inches(1.5), Inches(12), Inches(5), size=20, gap=10)

# ---------------- 5. Predictors ----------------
fig_slide("Predictors (50 features)", "feature_importance_by_group.png", side_bullets=[
    "Meteorology — ALADIN 2 km (T, wind, PBL, ventilation, radiation…).",
    "Proxy pollutants — PM₁₀, PM₂.₅, NO₂ (instrumented sites only).",
    "Terrain — DMR 3.5 (elevation, relief, multi-scale TPI, slope).",
    "Traffic — CDV, 2023 national census (load, HDV, dist. to road).",
    "Residential emission — bottom-up inventory (Krajčovičová 2020).",
    "No B[a]P autoregression (would leak / unavailable in gaps).",
], take="Static covariates make spatial generalization possible.")

# ---------------- 6. Model & validation ----------------
s = slide_blank(); header(s, "Model & validation design")
bullets(s, [
    "XGBoost: 1200 trees, lr 0.02, depth 6, mild L1/L2; sample weight 1/(B[a]P+0.5).",
    (1, "weighting trades a little pooled fit for lower rural bias (−17%) — see ablation."),
    "Skill = squared-Pearson r² (+ RMSE, MBE reported separately).",
    "Three cross-validation designs:",
    (1, "random k-fold — generic skill;"),
    (1, "leave-one-station-out (LOSO) — prediction at unmonitored sites;"),
    (1, "block-gap (30-day windows) — leakage-free gap-filling = primary metric."),
    "Baselines: persistence and seasonal climatology.",
], Inches(0.7), Inches(1.4), Inches(12), Inches(5.4), size=18, gap=7)

# ---------------- 7. Headline skill ----------------
fig_slide("Headline skill — block-gap vs baselines", "block_baseline_comparison.png",
    side_bullets=[
        "Block-gap median per-station r² = 0.78 (RMSE 1.51 ng m⁻³ pooled).",
        "≈ doubles seasonal climatology (r² 0.41).",
        "Far exceeds persistence (r² 0.23).",
        "Pooled block-gap r² (0.69) ≈ k-fold (0.70): sparse sampling ⇒ little leakage.",
    ], take="Reliable daily gap-filling at monitored stations.")

# ---------------- 8. Time series ----------------
fig_slide("Full-network gap-filled series", "timeseries_overview.png",
    take="Captures winter maxima / summer minima across station types.")

# ---------------- 9. Spatial generalization ----------------
fig_slide("Spatial generalization (LOSO) by station", "heatmap_stations.png",
    side_bullets=[
        "Generalizes well to urban regimes:",
        (1, "urban-bg r² ≈ 0.84, urban-traffic ≈ 0.83."),
        "Suburban background ≈ 0.75.",
        "Fails at the single industrial site (r² 0.05).",
        "Rural-background correlation weak (few analogues).",
    ], take="Hard cases are atypical, under-sampled regimes.")

# ---------------- 10. LOSO vs block-gap dissociation ----------------
fig_slide("Extending ≠ gap-filling", "barplot_r2.png", side_bullets=[
    "Atypical sites fail under LOSO but recover under block-gap:",
    (1, "Stará Lesná 0.00 → 0.68,"),
    (1, "SK0006R 0.38 → 0.59,"),
    (1, "industrial SK0018A 0.04 → 0.27."),
    "Withholding a whole atypical station removes its only analogue.",
    "Orographically distinct, high, remote sites are hardest to place.",
], take="Dependable gap-filler everywhere it is deployed; extension is the hard part.")

# ---------------- 11. Feature attribution ----------------
fig_slide("Feature attribution (TreeSHAP)", "shap_importance.png", side_bullets=[
    "By group: meteorology 32%, calendar/season 21%, proxies 19%,",
    (1, "bottom-up emission 11%, terrain 5%."),
    "Emission term is the 3rd-ranked individual predictor —",
    (1, "far above its split-gain rank."),
    "Terrain: small day-to-day, decisive cross-station.",
], take="Emission term does real predictive work.")

# ---------------- 12. Emission inventory result ----------------
s = slide_blank(); header(s, "Bottom-up emission inventory → rural bias fix", kicker="key novelty")
bullets(s, [
    "Adding the residential-emission term leaves aggregate skill ~unchanged…",
    (1, "at instrumented sites the proxies already encode emission intensity."),
    "…but it strongly corrects rural-background bias:",
    (1, "with proxies: RB mean |MBE| 0.40 → 0.20 ng m⁻³  (−50%),"),
    (1, "proxy-free (true virtual site): 0.81 → 0.25 ng m⁻³  (−69%)."),
    "Directly targets the largest uncertainty of existing European B[a]P maps (rural).",
    "Proxy-free virtual-site expectation: LOSO median r² ≈ 0.67.",
], Inches(0.7), Inches(1.4), Inches(12), Inches(5.4), size=19, gap=9)
takeaway(s, "The static emission field governs WHERE rural B[a]P sits.")

# ---------------- 13. Industrial regime ----------------
s = slide_blank(); header(s, "The industrial regime — what we tried", kicker="recent finding")
bullets(s, [
    "94% of national point-source B[a]P = U.S. Steel Košice coke ovens,",
    (1, "2 km from the only industrial station (Veľká Ida, SK0018A)."),
    "Physics is real: mean B[a]P 3.3 (wind away) → 7.9 ng m⁻³ (wind from plant);",
    (1, "partial corr 0.49 after controlling for dispersion + season."),
    "We built a wind-conditioned upwind industrial-emission term and retrained.",
    "Result: no cross-validated skill gain (block-gap 0.27→0.26, LOSO 0.05;",
    (1, "MBE only −2.24 → −2.07)."),
    "Cause is the network: a single industrial station leaves no analogue.",
], Inches(0.7), Inches(1.4), Inches(12), Inches(5.4), size=18, gap=7)
takeaway(s, "Limitation is network design — not a missing predictor.")

# ---------------- 14. Virtual network map ----------------
fig_slide("The virtual monitoring network", "virtual_network_map.png", side_bullets=[
    "Continuous daily B[a]P for all 54 cells, 2023–2025.",
    "Includes 33 never-monitored sites.",
    "17 of 54 sites exceed the 1 ng m⁻³ target (2024 annual mean).",
    "Highest: eastern industrial site & urban centres;",
    (1, "remote rural sites stay well below."),
], take="A usable continuous B[a]P product over Slovakia.")

# ---------------- 15. Limitations & next steps ----------------
s = slide_blank(); header(s, "Limitations & next steps")
bullets(s, [
    "Industrial regime: needs more industrial stations, not more emission data.",
    "Remote rural correlation stays modest (clean, near-detection-limit, sub-grid).",
    "Proxy pollutants exist only at instrumented sites → proxy-free bound r² ≈ 0.67.",
    "Next: continuous 2 km raster field; additional SI monitoring; longer record.",
], Inches(0.7), Inches(1.5), Inches(12), Inches(5), size=20, gap=12)

# ---------------- 16. Conclusions ----------------
s = slide_blank(); header(s, "Conclusions")
bullets(s, [
    "One globally trained model fills daily B[a]P gaps: block-gap median r² = 0.78.",
    "Clearly beats climatology and persistence.",
    "Terrain + traffic covariates are decisive for spatial generalization.",
    "Bottom-up residential emission cuts rural-background bias by up to 69%.",
    "Delivers a continuous virtual monitoring network, incl. 33 unmonitored sites.",
    "Open challenge — industrial regime — is a monitoring-network limitation.",
], Inches(0.7), Inches(1.5), Inches(12), Inches(5), size=20, gap=10)
takeaway(s, "Addresses the rural uncertainty that limits existing European B[a]P maps.")

out = os.path.join(OUTDIR, "bap_gapfilling_overview.pptx")
prs.save(out)
print(f"✅ {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")
